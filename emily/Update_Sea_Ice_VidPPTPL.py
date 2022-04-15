#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Created on Tue Apr 30 14:41:07 2019

@author: emily.niebuhr
"""

#!/usr/bin/env python3
# -*- coding: utf-8 -*-


#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Created on Fri Feb  8 10:22:54 2019

@author: emily.niebuhr
"""

import numpy as np

from docx import Document
from docx.shared import Inches
from pptx import Presentation


import urllib.request
import tarfile
from datetime import datetime, timedelta
from shutil import copy2
import os

import shutil 



from docx import Document
from docx.shared import Inches

##############################
## Webname
############################

RASM= str("RASM-ESRL_")
RASM_N= str("RASM-ESRL_4NIC_")

tg=str("-")
dg=str("-")
nc = str(".nc")
gzz = str(".tar.gz")


#####################
# DATE PL
#####################

x = datetime.now() 
x1=datetime.now() + timedelta(days=1)
x2=datetime.now() + timedelta(days=2)
x3=datetime.now() + timedelta(days=3)
x4=datetime.now() + timedelta(days=4)
x5=datetime.now() + timedelta(days=5)
x7=datetime.now() + timedelta(days=7)

mx1 = datetime.now() - timedelta(days=1)
mx2 = datetime.now() - timedelta(days=2)

mx5 = datetime.now() - timedelta(days=5)
mx8 = datetime.now() - timedelta(days=8)

mp1=x1.month
mp2=x2.month
mp3=x3.month
mp4=x4.month
mp5=x5.month
mp7=x7.month

yp1=x1.year
yp2=x2.year
yp3=x3.year
yp4=x4.year
yp5=x5.year
yp7=x7.year

mos = x.month
yr = x.year
daye = x.day
yester=mx1.day
ymos=mx1.month
yyear=mx1.year

yester22=mx2.day
ymos2=mx2.month
yyear2=mx2.year


yester5=mx5.day
y5mos=mx5.month
y5yr=mx5.year

###############################################3  Power Point ####################333


############## Change depending on wheather windows or linux machine 
os.path.dirname('/home/emily.niebuhr/Downloads')
dir1 = os.path.abspath('/home/emily.niebuhr/Downloads')
dir1 = os.path.abspath('/home/amos/IceOperatePL')
os.chdir(dir1)

name_EL = RASM_N + str(yyear)+tg+str("{:02d}".format(ymos))+tg+str("{:02d}".format(yester))+gzz 
dir_EL = RASM_N + str(yyear)+tg+str("{:02d}".format(ymos))+tg+str("{:02d}".format(yester)) 
dir2_EL = RASM+str(yyear)+tg+str("{:02d}".format(ymos))+tg+str("{:02d}".format(yester)) 

#################################33
os.chdir(dir1+"/"+dir_EL+"/"+dir2_EL)
print (os.getcwd())

prs = Presentation()

title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Sea Ice Videos  " + str("{:02d}".format(mos))+"/"+str("{:02d}".format(daye))

#########  Slide Two ##################
title_content_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(title_content_layout)
title = slide.shapes.title

#title.text = "5 Day Ice Forecasts."
top = Inches(0.5)  # how many inches from top of slide
top2 = Inches(3.8)  # how many inches from top of slide

left =  Inches(0.5)  # one inch from left part of slide
height = Inches(6.5) # height in inches of image
width = Inches(8.5)  # width in inches of image
## Ok lets add this image
top = Inches(0.2)  # how many inches from top of slide
pic = slide.shapes.add_picture("GIOPSV5.gif", left, top, height=height,width=width)
#########  Slide Two ##################


title_content_layout = prs.slide_layouts[2]
slide = prs.slides.add_slide(title_content_layout)
title = slide.shapes.title

#title.text = "5 Day Ice Forecasts."

######## Image Download Time ######################
##  Note- to determine image properties, open your ppt, right click on image
## and choose size and position.  Then use these numbers to set your image size and 
## location
#  width=Inches(3.1), height=Inches(2.2)
#top = Inches(1.5)  # how many inches from top of slide
top = Inches(0.5)  # how many inches from top of slide
top2 = Inches(3.8)  # how many inches from top of slide

left =  Inches(0.5)  # one inch from left part of slide
height = Inches(6.5) # height in inches of image
width = Inches(8.5)  # width in inches of image
## Ok lets add this image
top = Inches(0.2)  # how many inches from top of slide
pic = slide.shapes.add_picture("ESRL10.gif", left, top, height=height,width=width)



################################
#########  Slide Two ##################


title_content_layout = prs.slide_layouts[3]
slide = prs.slides.add_slide(title_content_layout)
title = slide.shapes.title

#title.text = "5 Day Ice Forecasts."

top = Inches(0.5)  # how many inches from top of slide
top2 = Inches(3.8)  # how many inches from top of slide

left =  Inches(0.5)  # one inch from left part of slide
height = Inches(6.5) # height in inches of image
width = Inches(8.5)  # width in inches of image
## Ok lets add this image
top = Inches(0.2)  # how many inches from top of slide
pic = slide.shapes.add_picture("GOFS_goLoopPL.gif", left, top, height=height,width=width)



################################
#########  Slide Two ##################


title_content_layout = prs.slide_layouts[3]
slide = prs.slides.add_slide(title_content_layout)
title = slide.shapes.title

#title.text = "5 Day Ice Forecasts."

top = Inches(0.5)  # how many inches from top of slide
top2 = Inches(3.8)  # how many inches from top of slide

left =  Inches(0.5)  # one inch from left part of slide
height = Inches(6.5) # height in inches of image
width = Inches(8.5)  # width in inches of image
## Ok lets add this image
top = Inches(0.2)  # how many inches from top of slide
pic = slide.shapes.add_picture("RTOFS_Loop.gif", left, top, height=height,width=width)



#########  Slide Two ##################


title_content_layout = prs.slide_layouts[4]
slide = prs.slides.add_slide(title_content_layout)
title = slide.shapes.title

top = Inches(0.5)  # how many inches from top of slide
top2 = Inches(3.8)  # how many inches from top of slide

left =  Inches(0.5)  # one inch from left part of slide
height = Inches(6.5) # height in inches of image
width = Inches(8.5)  # width in inches of image
## Ok lets add this image
top = Inches(0.2)  # how many inches from top of slide
pic = slide.shapes.add_picture("TH_GOFS_goLoopPL.gif", left, top, height=height,width=width)

title_content_layout = prs.slide_layouts[4]
slide = prs.slides.add_slide(title_content_layout)
title = slide.shapes.title

#title.text = "5 Day Ice Forecasts."

######## Image Download Time ######################

top = Inches(0.5)  # how many inches from top of slide
top2 = Inches(3.8)  # how many inches from top of slide

left =  Inches(0.5)  # one inch from left part of slide
height = Inches(6.5) # height in inches of image
width = Inches(8.5)  # width in inches of image
## Ok lets add this image
top = Inches(0.2)  # how many inches from top of slide
pic = slide.shapes.add_picture("RTOFS_thLoop.gif", left, top, height=height,width=width)

######################
#########  Slide Two ##################


title_content_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_content_layout)
title = slide.shapes.title

#title.text = "5 Day Ice Forecasts."

######## Image Download Time ######################
#url= "https://www.weather.gov/images/afc/ice/CT.jpg"
#names =  urllib.request.urlretrieve(url,"ice_puppy1.png")

##  Note- to determine image properties, open your ppt, right click on image
## and choose size and position.  Then use these numbers to set your image size and 
## location
#  width=Inches(3.1), height=Inches(2.2)
#top = Inches(1.5)  # how many inches from top of slide
top = Inches(0.5)  # how many inches from top of slide
top2 = Inches(3.8)  # how many inches from top of slide

left =  Inches(0.5)  # one inch from left part of slide
height = Inches(6.5) # height in inches of image
width = Inches(8.5)  # width in inches of image
## Ok lets add this image
top = Inches(0.2)  # how many inches from top of slide
pic = slide.shapes.add_picture("TH_GIOPSV5.gif", left, top, height=height,width=width)

tg=str("_")





title_content_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(title_content_layout)
title = slide.shapes.title

#title.text = "5 Day Ice Forecasts."

######## Image Download Time ######################
#url= "https://www.weather.gov/images/afc/ice/CT.jpg"
#names =  urllib.request.urlretrieve(url,"ice_puppy1.png")

##  Note- to determine image properties, open your ppt, right click on image
## and choose size and position.  Then use these numbers to set your image size and 
## location
#  width=Inches(3.1), height=Inches(2.2)
#top = Inches(1.5)  # how many inches from top of slide
top = Inches(0.5)  # how many inches from top of slide
top2 = Inches(3.8)  # how many inches from top of slide

left =  Inches(0.5)  # one inch from left part of slide
height = Inches(6.5) # height in inches of image
width = Inches(8.5)  # width in inches of image
## Ok lets add this image
top = Inches(0.2)  # how many inches from top of slide
pic = slide.shapes.add_picture("TH_ESRL10.gif", left, top, height=height,width=width)


title_content_layout = prs.slide_layouts[7]
slide = prs.slides.add_slide(title_content_layout)
title = slide.shapes.title

#title.text = "5 Day Ice Forecasts."

######## Image Download Time ######################
#url= "https://www.weather.gov/images/afc/ice/CT.jpg"
#names =  urllib.request.urlretrieve(url,"ice_puppy1.png")

##  Note- to determine image properties, open your ppt, right click on image
## and choose size and position.  Then use these numbers to set your image size and 
## location
#  width=Inches(3.1), height=Inches(2.2)
#top = Inches(1.5)  # how many inches from top of slide
###top = Inches(4.5)  # how many inches from top of slide
###top2 = Inches(4.5)  # how many inches from top of slide
#1.5
###left =  Inches(0.5)  # one inch from left part of slide
###height = Inches(4.0) # height in inches of image
###width = Inches(3.5)  # width in inches of image
## Ok lets add this image
###top = Inches(3.5)  # how many inches from top of slide
###pic = slide.shapes.add_picture("Synoptic_Today1.gif", left, top, height=height,width=width)

##### download and add second image.
#url= "https://www.weather.gov/images/afc/ice/Forecast.jpg"
#names2 =  urllib.request.urlretrieve(url,"ice_puppy2.png")

###left  = Inches(5.0)  # lets put this image 6 inches from left hand side
###pic = slide.shapes.add_picture("Synoptic_FiveDay.gif", left, top, height=height,width=width)
#left =  Inches(0.5)  # one inch from left part of slide

#pic = slide.shapes.add_picture("plot_Day_GIOPS1.png", left, top2, height=height,width=width)
#left  = Inches(5)  # lets put this image 6 inches from left hand side

#pic = slide.shapes.add_picture("plot_hour12_RIOPS1.png", left, top2, height=height,width=width)


top = Inches(0.5)  # how many inches from top of slide
top2 = Inches(3.8)  # how many inches from top of slide

left =  Inches(0.5)  # one inch from left part of slide
height = Inches(6.5) # height in inches of image
width = Inches(8.5)  # width in inches of image
## Ok lets add this image
top = Inches(0.2)  # how many inches from top of slide
pic = slide.shapes.add_picture("Syn_Loop.gif", left, top, height=height,width=width)



title_content_layout = prs.slide_layouts[8]
slide = prs.slides.add_slide(title_content_layout)
title = slide.shapes.title

top = Inches(4.5)  # how many inches from top of slide
top2 = Inches(4.5)  # how many inches from top of slide
#1.5
left =  Inches(0.5)  # one inch from left part of slide
height = Inches(4.0) # height in inches of image
width = Inches(3.5)  # width in inches of image
## Ok lets add this image
top = Inches(3.5)  # how many inches from top of slide
pic = slide.shapes.add_picture("ASIP_currConc.jpg", left, top, height=height,width=width)


left  = Inches(5.0)  # lets put this image 6 inches from left hand side
pic = slide.shapes.add_picture("curr_full_stage.jpg", left, top, height=height,width=width)
tg=str("_")


title_content_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_content_layout)
title = slide.shapes.title

#title.text = "5 Day Ice Forecasts."

######## Image Download Time ######################

top = Inches(0.5)  # how many inches from top of slide
top2 = Inches(3.8)  # how many inches from top of slide

left =  Inches(0.5)  # one inch from left part of slide
height = Inches(6.5) # height in inches of image
width = Inches(8.5)  # width in inches of image
## Ok lets add this image
top = Inches(0.2)  # how many inches from top of slide
pic = slide.shapes.add_picture("WW3_Waves.gif", left, top, height=height,width=width)



title_content_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_content_layout)
title = slide.shapes.title
######## Image Download Time ######################

top = Inches(0.5)  # how many inches from top of slide
top2 = Inches(3.8)  # how many inches from top of slide

left =  Inches(0.5)  # one inch from left part of slide
height = Inches(6.5) # height in inches of image
width = Inches(8.5)  # width in inches of image
## Ok lets add this image
top = Inches(0.2)  # how many inches from top of slide
pic = slide.shapes.add_picture("WWW3_Winds.gif", left, top, height=height,width=width)





datename = "Sea_Ice_Videos_"+str("{:02d}".format(mos))+tg+str("{:02d}".format(daye))+str("NWPS_Video_PPT.pptx")
prs.save(datename)

